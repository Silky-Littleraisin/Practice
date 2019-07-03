# -*- coding: utf-8 -*-
from pptx import Presentation

prs = Presentation("basic.pptx")

text_runs = []
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

print("\n".join(text_runs))

for p, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            text_runs = [run.text for run in paragraph.runs]
            print(
                p,
                shape.placeholder_format.type,
                paragraph.level,
                " ".join(text_runs),
                sep="\t")