# -*- coding: utf-8 -*-
from pptx import Presentation

# スライドオブジェクト
prs = Presentation()
# スライドレイアウトの指定
# 0, Title (presentation title slide)
# 1, Title and Content
# 2, Section Header (sometimes called Segue)
# 3, Two Content (side by side bullet textboxes)
# 4, Comparison (same but additional title for each side by side content box)
# 5, Title Only
# 6, Blank
# 7, Content with Caption
# 8, Picture with Caption
title_slide_layout = prs.slide_layouts[0]

# タイトルスライド作成
slide = prs.slides.add_slide(title_slide_layout)
# タイトル
slide.shapes.title.text = "デジタルドキュメント論"
# サブタイトル
slide.placeholders[1].text = "東京大学教育学研究科"

tf = slide.shapes.placeholders[1].text_frame
p = tf.add_paragraph()  # 段落を追加
p.text = "2019年度 S1, S2"
p.level = 1 # 段落のレベル

# 標準スライドレイアウト作成
bullet_slide_layout = prs.slide_layouts[1]
# スライドページの追加
slide = prs.slides.add_slide(bullet_slide_layout)

shapes = slide.shapes
shapes.title.text = '本日の内容'

tf = shapes.placeholders[1].text_frame
tf.text = 'Word形式'

p = tf.add_paragraph()
p.text = 'Word形式のファイルの作成'
p.level = 1

p = tf.add_paragraph()
p.text = 'python-docxを使ってdocxファイルを作成する'
p.level = 2

contents = [
    ['PowerPoint形式', 0],
    ['PowerPoint形式のファイルの作成', 1],
    ['python-pptxを使ってpptxファイルを作成する', 2]
]
for row in contents:
    p = tf.add_paragraph()
    p.text = row[0]
    p.level = row[1]

prs.save('basic.pptx')