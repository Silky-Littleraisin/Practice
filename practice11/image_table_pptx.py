#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches

# 既存のpptxファイルを読み込む
prs = Presentation("basic.pptx")

slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = '画像の配置'

# 画像
height = Inches(3.0)
top = Inches(2.0)
left = Inches(0.75)
pic = slide.shapes.add_picture('wordcloud3.png', left, top, height=height)

left = Inches(5.25)
pic = slide.shapes.add_picture('wordcloud4.png', left, top, height=height)

# テキストボックス
width = height = Inches(1)
left = Inches(1.0)
top = Inches(5.25)
txBox = slide.shapes.add_textbox(left, top, width, height)
txBox.text_frame.text = "ワードクラウド3"

left = Inches(5.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
txBox.text_frame.text = "ワードクラウド4"


# 2ページ目
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
shapes.title.text = '表の配置'

rows = 4  # 行数
cols = 2  # 列数
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(1.0)
table = shapes.add_table(rows, cols, left, top, width, height).table

# 列の長さ
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# 列の見出し
table.cell(0, 0).text = '日にち'
table.cell(0, 1).text = '内容'

# 行の本文
table.cell(1, 0).text = '4/10'
table.cell(1, 1).text = 'ガイダンス'

table.cell(2, 0).text = '4/17'
table.cell(2, 1).text = '学術論文PDF解析'

table.cell(3, 0).text = '4/24'
table.cell(3, 1).text = '小説のテキストマイニング'

prs.save('image_table.pptx')